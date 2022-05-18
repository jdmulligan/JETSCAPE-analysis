"""
  macro for plotting analyzed jetscape events
  """

# This script plots histograms created in the analysis of Jetscape events
#
# Author: James Mulligan (james.mulligan@berkeley.edu)

# General
import ctypes
import os
import sys
import yaml
import argparse

# Data analysis and plotting
import ROOT
import numpy as np
import uproot
import h5py

# Base class
sys.path.append('.')
from jetscape_analysis.base import common_base
from plot import plot_results_STAT_utils

# Prevent ROOT from stealing focus when plotting
ROOT.gROOT.SetBatch(True)

################################################################
class PlotResults(common_base.CommonBase):

    # ---------------------------------------------------------------
    # Constructor
    # ---------------------------------------------------------------
    def __init__(self, config_file='', input_file='', output_dir='', pp_ref_file='', **kwargs):
        super(PlotResults, self).__init__(**kwargs)

        if output_dir:
            self.output_dir = output_dir
            if not os.path.exists(self.output_dir):
                os.makedirs(self.output_dir)
        else:
            self.output_dir = os.path.dirname(input_file)

        self.plot_utils = plot_results_STAT_utils.PlotUtils()
        self.plot_utils.setOptions()
        ROOT.gROOT.ForceStyle()

        self.input_file = ROOT.TFile(input_file, 'READ')

        self.data_color = ROOT.kGray+3
        self.data_marker = 21
        self.jetscape_color = [ROOT.kViolet-8, ROOT.kViolet-8, ROOT.kRed-7, ROOT.kTeal-8, ROOT.kCyan-2, ROOT.kGreen-6, ROOT.kAzure-4, ROOT.kOrange+6, ROOT.kBlue-10]
        self.jetscape_fillstyle = [1001, 3144, 1001, 3144]
        self.jetscape_alpha = [0.7, 0.7, 0.7, 0.7]
        self.jetscape_marker = 20
        self.marker_size = 1.5
        self.line_width = 2
        self.line_style = 1
        self.file_format = '.pdf'

        # Check whether pp or AA
        if 'PbPb' in input_file or 'AuAu' in input_file:
            self.is_AA = True
            self.observable_centrality_list = []
        else:
            self.is_AA = False

        # If AA, load the pp reference results so that we can construct RAA
        if self.is_AA:
            self.pp_ref_file = ROOT.TFile(pp_ref_file, 'READ')

        # Read config file
        with open(config_file, 'r') as stream:
            self.config = yaml.safe_load(stream)
        self.sqrts = self.config['sqrt_s']
        self.power = self.config['power']
        self.pt_ref = self.config['pt_ref']

        # If AA, set different options for hole subtraction treatment
        self.jet_collection_labels_AA = self.config['jet_collection_labels'] + ['shower_recoil_unsubtracted']
        self.jet_collection_label_pp = ''
        if self.is_AA:
            self.jet_collection_labels = self.jet_collection_labels_AA
        else:
            self.jet_collection_labels = [self.jet_collection_label_pp]

        # We will write final results after all scalings, along with data, to file
        self.output_dict = {}

        print(self)

    #-------------------------------------------------------------------------------------------
    #-------------------------------------------------------------------------------------------
    #-------------------------------------------------------------------------------------------
    def plot_results(self):

        self.analysis = 'Analysis1'
        
        if self.analysis == 'Analysis1':

            self.plot_hadron_observables(observable_type='hadron')
                        
            self.plot_jet_observables(observable_type='inclusive_chjet')

            if 'inclusive_jet' in self.config:
                self.plot_jet_observables(observable_type='inclusive_jet')

        else:

            self.plot_hadron_observables(observable_type='hadron')

            self.plot_hadron_correlation_observables(observable_type='hadron_correlations')

            self.plot_jet_observables(observable_type='inclusive_chjet')

            if 'inclusive_jet' in self.config:
                self.plot_jet_observables(observable_type='inclusive_jet')

            if 'semi_inclusive_chjet' in self.config:
                self.plot_semi_inclusive_chjet_observables(observable_type='semi_inclusive_chjet')

            if 'dijet' in self.config:
                self.plot_jet_observables(observable_type='dijet')

        self.plot_event_qa()

        self.write_output_objects()

        # Generate pptx for convenience
        if self.file_format == '.png':
            self.generate_pptx()

    #-------------------------------------------------------------------------------------------
    # Plot hadron observables
    #-------------------------------------------------------------------------------------------
    def plot_hadron_observables(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')

        for observable, block in self.config[observable_type].items():
            for centrality_index,centrality in enumerate(block['centrality']):

                if 'hepdata' not in block and 'custom_data' not in block:
                    continue

                # Initialize observable configuration
                self.suffix = ''
                self.init_observable(observable_type, observable, block, centrality, centrality_index)

                # Plot observable
                self.plot_observable(observable_type, observable, centrality)

    #-------------------------------------------------------------------------------------------
    # Plot hadron correlation observables
    #-------------------------------------------------------------------------------------------
    def plot_hadron_correlation_observables(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')

        for observable, block in self.config[observable_type].items():
            for centrality_index, centrality in enumerate(block['centrality']):

                if 'hepdata' not in block and 'custom_data' not in block:
                    continue

                #for hadron v2
                if 'v2' in observable:
                    self.init_observable(observable_type, observable, block, centrality, centrality_index)
                    # Histogram observable
                    self.plot_observable(observable_type, observable, centrality)

                # STAR dihadron
                if observable == 'dihadron_star':

                    # Initialize observable configuration
                    pt_trigger_ranges = block["pt_trig"]
                    pt_associated_ranges = block["pt_assoc"]
                    # Loop over trigger and associated ranges
                    for pt_trig_min, pt_trig_max in pt_trigger_ranges:
                        for pt_assoc_min, pt_assoc_max in pt_associated_ranges:
                            # If the upper range has -1, it's unbounded, so we make it large enough not to matter
                            pt_assoc_max = 1000 if pt_assoc_max == -1 else pt_assoc_max
                            self.suffix = f"_pt_trig_{pt_trig_min:g}_{pt_trig_max:g}_pt_assoc_{pt_assoc_min:g}_{pt_assoc_max:g}"
                            self.init_observable(observable_type, observable, block, centrality, centrality_index)

                            # Histogram observable
                            self.plot_observable(observable_type, observable, centrality)

    #-------------------------------------------------------------------------------------------
    # Histogram inclusive jet observables
    #-------------------------------------------------------------------------------------------
    def plot_jet_observables(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')

        for observable, block in self.config[observable_type].items():
            for centrality_index,centrality in enumerate(block['centrality']):

                if self.analysis == 'Analysis1':
                    if not observable.startswith('pt_'):
                        print(f'skipping {observable}')
                        continue

                for self.jet_R in block['jet_R']:

                    # Optional: Loop through pt bins
                    for pt_bin in range(len(block['pt'])-1):

                        if len(block['pt']) > 2:
                            pt_suffix = f'_pt{pt_bin}'
                        else:
                            pt_suffix = ''

                        # Optional: subobservable
                        subobservable_label_list = ['']
                        if 'kappa' in block:
                            subobservable_label_list = [f'_k{kappa}' for kappa in block['kappa']]
                        for subobservable_label in subobservable_label_list:

                            # Set normalization
                            self_normalize = False
                            for x in ['mass', 'g', 'ptd', 'charge', 'mg', 'zg', 'tg', 'xj']:
                                if x in observable:
                                    self_normalize = True

                            if 'SoftDrop' in block:
                                for grooming_setting in block['SoftDrop']:
                                    if observable == 'zg_alice' or observable == 'tg_alice':
                                        if np.isclose(self.jet_R, 0.4) and centrality_index == 0:
                                            continue
                                        if np.isclose(self.jet_R, 0.2) and centrality_index == 1:
                                            continue

                                    print(f'      grooming_setting = {grooming_setting}')
                                    zcut = grooming_setting['zcut']
                                    beta = grooming_setting['beta']

                                    self.suffix = f'_R{self.jet_R}_zcut{zcut}_beta{beta}{subobservable_label}'
                                    if 'hepdata' not in block and 'custom_data' not in block:
                                        continue

                                    # Initialize observable configuration
                                    self.init_observable(observable_type, observable, block, centrality, centrality_index, pt_suffix=pt_suffix, self_normalize=self_normalize)

                                    # Plot observable
                                    self.plot_observable(observable_type, observable, centrality, pt_suffix)

                            else:

                                self.suffix = f'_R{self.jet_R}{subobservable_label}'
                                if 'hepdata' not in block and 'custom_data' not in block:
                                    continue

                                # Initialize observable configuration
                                self.init_observable(observable_type, observable, block, centrality, centrality_index, pt_suffix=pt_suffix, self_normalize=self_normalize)

                                # Plot observable
                                self.plot_observable(observable_type, observable, centrality, pt_suffix)

    #-------------------------------------------------------------------------------------------
    # Histogram semi-inclusive jet observables
    #-------------------------------------------------------------------------------------------
    def plot_semi_inclusive_chjet_observables(self, observable_type=''):
        print()
        print(f'Plot {observable_type} observables...')

        for observable, block in self.config[observable_type].items():
            for centrality_index,centrality in enumerate(block['centrality']):

                for self.jet_R in block['jet_R']:

                    self.suffix = f'_R{self.jet_R}'
                    if 'hepdata' not in block and 'custom_data' not in block:
                        continue

                    # Set normalization
                    self_normalize = False
                    for x in ['nsubjettiness']:
                        if x in observable:
                            self_normalize = True

                    # Initialize observable configuration
                    self.init_observable(observable_type, observable, block, centrality, centrality_index, self_normalize=self_normalize)

                    # Plot observable
                    self.plot_observable(observable_type, observable, centrality)

    #-------------------------------------------------------------------------------------------
    # Initialize a single observable's config
    #-------------------------------------------------------------------------------------------
    def init_observable(self, observable_type, observable, block, centrality, centrality_index, pt_suffix='', self_normalize=False):

        # Initialize an empty dict containing relevant info
        self.observable_settings = {}

        # Initialize common settings into class members
        self.init_common_settings(observable, block)

        #-----------------------------------------------------------
        # Initialize data distribution into self.observable_settings
        if 'hepdata' in block:
            self.observable_settings['data_distribution'] = self.plot_utils.tgraph_from_hepdata(block, self.is_AA, self.sqrts, observable_type, observable, centrality_index, suffix=self.suffix, pt_suffix=pt_suffix)
        elif 'custom_data' in block:
            self.observable_settings['data_distribution'] = self.plot_utils.tgraph_from_yaml(block, self.is_AA, self.sqrts, observable_type, observable, centrality_index, suffix=self.suffix, pt_suffix=pt_suffix)
        else:
            self.observable_settings['data_distribution'] = None

        #-----------------------------------------------------------
        # Initialize JETSCAPE distribution into self.observable_settings
        self.initialize_jetscape_distribution(observable_type, observable, block, centrality, centrality_index, pt_suffix=pt_suffix, self_normalize=self_normalize)

        #-----------------------------------------------------------
        # For pp case -- form ratio of JETSCAPE to data and load into self.observable_settings
        if not self.is_AA:
            if self.observable_settings['data_distribution'] and self.observable_settings[f'jetscape_distribution'] and not self.observable_settings[f'jetscape_distribution'].InheritsFrom(ROOT.TH2.Class()) and not self.skip_pp_ratio:
                self.observable_settings['ratio'] = self.plot_utils.divide_histogram_by_tgraph(self.observable_settings[f'jetscape_distribution'],
                                                                                self.observable_settings['data_distribution'], include_tgraph_uncertainties=False)
            else:
                self.observable_settings['ratio'] = None
        if 'v2' in observable and self.is_AA and self.observable_settings[f'jetscape_distribution']:
            self.observable_settings['ratio'] =  self.plot_utils.divide_histogram_by_tgraph(self.observable_settings[f'jetscape_distribution'], self.observable_settings['data_distribution'])

    #-------------------------------------------------------------------------------------------
    # Initialize from settings from config file into class members
    #-------------------------------------------------------------------------------------------
    def init_common_settings(self, observable, block):

        self.xtitle = block['xtitle']
        if 'eta_cut' in block:
            self.eta_cut = block['eta_cut']
        if 'y_cut' in block:
            self.y_cut = block['y_cut']
        if 'pt' in block:
            self.pt = block['pt']
        if 'eta_cut_R' in block:
            self.eta_R = block['eta_cut_R']
            self.eta_cut = np.round(self.eta_R - self.jet_R, decimals=1)
        if 'c_ref' in block:
            index = block['jet_R'].index(self.jet_R)
            self.c_ref = block['c_ref'][index]
        if 'low_trigger_range' in block:
            self.low_trigger_range = block['low_trigger_range']
        if 'high_trigger_range' in block:
            self.high_trigger_range = block['high_trigger_range']
        if 'trigger_range' in block:
            self.trigger_range = block['trigger_range']
        if 'logy' in block:
            self.logy = block['logy']
        else:
            self.logy = False

        #for v2
        if 'v2' in observable:
            self.y_ratio_min = -0.5
            self.y_ratio_max = 1.99

        if self.is_AA:
            if 'ytitle_AA' in block:
                self.ytitle = block['ytitle_AA']
            if 'y_min_AA' in block:
                self.y_min = float(block['y_min_AA'])
                self.y_max = float(block['y_max_AA'])
            else:
                self.y_min = 0.
                self.y_max = 1.
        else:
            if 'ytitle_pp' in block:
                self.ytitle = block['ytitle_pp']
            else:
                self.ytitle = ''
            if 'y_min_pp' in block:
                self.y_min = float(block['y_min_pp'])
                self.y_max = float(block['y_max_pp'])
            else:
                self.y_min = 0.
                self.y_max = 1.99
            if 'y_ratio_min' in block:
                self.y_ratio_min = block['y_ratio_min']
                self.y_ratio_max = block['y_ratio_max']
            else:
                self.y_ratio_min = 0.
                self.y_ratio_max = 1.99
            # Provide the opportunity to disable logy in pp, but keep in AA
            if 'logy_pp' in block:
                self.logy = block["logy_pp"]

        if 'skip_pp' in block:
            self.skip_pp = block['skip_pp']
        else:
            self.skip_pp = False
        if 'skip_pp_ratio' in block:
            self.skip_pp_ratio = block['skip_pp_ratio']
        else:
            self.skip_pp_ratio = False
        if 'skip_AA_ratio' in block:
            self.skip_AA_ratio = block['skip_AA_ratio']
        else:
            self.skip_AA_ratio = False
        if 'scale_by' in block:
            self.scale_by = block['scale_by']
        else:
            self.scale_by = None

        # Flag to plot hole histogram (for hadron histograms only)
        if self.is_AA:
            self.subtract_holes = observable in ['pt_ch_alice', 'pt_pi_alice', 'pt_pi0_alice', 'pt_ch_cms',
                                                 'pt_ch_atlas', 'pt_pi0_phenix', 'pt_ch_star', 'v2_atlas', 'v2_cms',
                                                 'dihadron_star']
        else:
            self.subtract_holes = False

    #-------------------------------------------------------------------------------------------
    # Initialize JETSCAPE distribution into self.observable_settings
    #-------------------------------------------------------------------------------------------
    def initialize_jetscape_distribution(self, observable_type, observable, block, centrality, centrality_index, pt_suffix='', self_normalize=False):

        #-------------------------------------------------------------
        # AA
        if self.is_AA:

            # Add centrality bin to list, if needed
            if centrality not in self.observable_centrality_list:
                self.observable_centrality_list.append(centrality)

            #-------------------------------------------------------------
            # For hadron observables, we retrieve and subtract the hole histograms
            if self.subtract_holes:

                hole_labels = ['', '_holes']
                for hole_label in hole_labels:

                    # Get histogram, and add to self.observable_settings
                    self.get_histogram(observable_type, observable, centrality, collection_label=hole_label, pt_suffix=pt_suffix)

                    # Normalization
                    # Note: If we divide by the sum of weights (corresponding to n_events) and multiply by the
                    #       pt-hat cross-section, then JETSCAPE distribution gives cross-section: dsigma/dx (in mb)
                    self.scale_histogram(observable_type, observable, centrality,
                                         collection_label=hole_label, pt_suffix=pt_suffix, self_normalize=self_normalize)

                # Subtract the holes (and save unsubtracted histogram)
                if self.observable_settings[f'jetscape_distribution']:
                    self.observable_settings['jetscape_distribution_unsubtracted'] = self.observable_settings[f'jetscape_distribution'].Clone()
                    self.observable_settings['jetscape_distribution_unsubtracted'].SetName('{}_unsubtracted'.format(self.observable_settings[f'jetscape_distribution'].GetName()))
                    if self.observable_settings['jetscape_distribution_holes']:
                        self.observable_settings['jetscape_distribution'].Add(self.observable_settings['jetscape_distribution_holes'], -1)

                # Perform any additional manipulations on scaled histograms
                self.post_process_histogram(observable_type, observable, block, centrality, centrality_index)


            #-------------------------------------------------------------
            # For jet histograms, loop through all available hole subtraction variations, and initialize histogram
            else:

                for jet_collection_label in self.jet_collection_labels:

                    # Get histogram and add to self.observable_settings
                    #  - In the case of semi-inclusive measurements construct difference of histograms
                    self.get_histogram(observable_type, observable, centrality, pt_suffix=pt_suffix, collection_label=jet_collection_label)
                    self.scale_histogram(observable_type, observable, centrality,
                                        collection_label=jet_collection_label, pt_suffix=pt_suffix, self_normalize=self_normalize)
                    self.post_process_histogram(observable_type, observable, block, centrality, centrality_index, collection_label=jet_collection_label)

        #-------------------------------------------------------------
        # pp
        else:
            self.get_histogram(observable_type, observable, centrality, pt_suffix=pt_suffix)
            self.scale_histogram(observable_type, observable, centrality, pt_suffix=pt_suffix, self_normalize=self_normalize)
            self.post_process_histogram(observable_type, observable, block, centrality, centrality_index)
    #-------------------------------------------------------------------------------------------
    # Get histogram and add to self.observable_settings
    #  - In AA case, also add hole histogram
    #  - In the case of semi-inclusive measurements construct difference of histograms
    #-------------------------------------------------------------------------------------------
    def get_histogram(self, observable_type, observable, centrality, collection_label='', pt_suffix=''):

        keys = [key.ReadObj().GetTitle() for key in self.input_file.GetListOfKeys()]

        # In the case of semi-inclusive measurements construct difference of histograms
        if 'semi_inclusive' in observable_type:
            self.construct_semi_inclusive_histogram(keys, observable_type, observable, centrality, collection_label=collection_label)

        # For all other histograms, get the histogram directly
        else:

            # Get histogram
            self.hname = f'h_{observable_type}_{observable}{self.suffix}{collection_label}_{centrality}{pt_suffix}'
            if self.hname in keys:
                h_jetscape = self.input_file.Get(self.hname)
                h_jetscape.SetDirectory(0)
                if not h_jetscape.GetSumw2():
                    h_jetscape.Sumw2()

            else:
                h_jetscape = None
            self.observable_settings[f'jetscape_distribution{collection_label}'] = h_jetscape

    #-------------------------------------------------------------------------------------------
    # Construct semi-inclusive observables from difference of histograms
    #-------------------------------------------------------------------------------------------
    def construct_semi_inclusive_histogram(self, keys, observable_type, observable, centrality, collection_label=''):

        if self.sqrts == 2760: # Delta recoil

            hname_low_trigger = f'h_{observable_type}_{observable}_R{self.jet_R}_lowTrigger{collection_label}_{centrality}'
            hname_high_trigger = f'h_{observable_type}_{observable}_R{self.jet_R}_highTrigger{collection_label}_{centrality}'
            hname_ntrigger = f'h_{observable_type}_alice_trigger_pt{observable}{collection_label}_{centrality}'
            self.hname = f'h_{observable_type}_{observable}_R{self.jet_R}{collection_label}_{centrality}'
            if hname_low_trigger in keys and hname_high_trigger in keys and hname_ntrigger in keys:
                h_jetscape_low = self.input_file.Get(hname_low_trigger)
                h_jetscape_low.SetDirectory(0)
                h_jetscape_high = self.input_file.Get(hname_high_trigger)
                h_jetscape_high.SetDirectory(0)
                h_jetscape_ntrigger = self.input_file.Get(hname_ntrigger)
                h_jetscape_ntrigger.SetDirectory(0)

                low_trigger = (self.low_trigger_range[0]+self.low_trigger_range[1])/2
                high_trigger = (self.high_trigger_range[0]+self.high_trigger_range[1])/2
                n_trig_high = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(high_trigger))
                n_trig_low = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(low_trigger))
                h_jetscape_high.Scale(1./n_trig_high)
                h_jetscape_low.Scale(1./n_trig_low)
                self.observable_settings[f'jetscape_distribution{collection_label}'] = h_jetscape_high.Clone('h_delta_recoil')
                self.observable_settings[f'jetscape_distribution{collection_label}'].Add(h_jetscape_low, -1)
            else:
                self.observable_settings[f'jetscape_distribution{collection_label}'] = None

        elif self.sqrts == 200:

            hname = f'h_{observable_type}_{observable}_R{self.jet_R}{collection_label}_{centrality}'
            hname_ntrigger = f'h_{observable_type}_star_trigger_pt{collection_label}_{centrality}'
            self.hname = f'h_{observable_type}_{observable}_R{self.jet_R}{collection_label}_{centrality}'
            if hname in keys and hname_ntrigger in keys:
                self.observable_settings[f'jetscape_distribution{collection_label}'] = self.input_file.Get(hname)
                self.observable_settings[f'jetscape_distribution{collection_label}'].SetDirectory(0)
                h_jetscape_ntrigger = self.input_file.Get(hname_ntrigger)
                h_jetscape_ntrigger.SetDirectory(0)

                trigger = (self.trigger_range[0]+self.trigger_range[1])/2
                n_trig = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(trigger))
                self.observable_settings[f'jetscape_distribution{collection_label}'].Scale(1./n_trig)
            else:
                self.observable_settings[f'jetscape_distribution{collection_label}'] = None

    #-------------------------------------------------------------------------------------------
    # Scale histogram for:
    #  - xsec/weight_sum
    #  - bin width
    #  - observable-specific factors (eta, sigma_inel, n_jets, ...)
    #  - self-normalization
    #-------------------------------------------------------------------------------------------
    def scale_histogram(self, observable_type, observable, centrality, collection_label='', pt_suffix='', self_normalize=False):

        h = self.observable_settings[f'jetscape_distribution{collection_label}']
        if not h:
            return

        # (0) No scaling is needed for hadron v2 and jet v2
        if 'v2' in observable:
            return

        #--------------------------------------------------
        # (1) Scale all histograms by the min-pt-hat cross-section and weight-sum
        if self.is_AA:
            h_xsec = self.input_file.Get(f'h_xsec_{centrality}')
            weight_sum = self.input_file.Get(f'h_weight_sum_{centrality}').GetBinContent(1)
        else:
            h_xsec = self.input_file.Get('h_xsec')
            weight_sum = self.input_file.Get('h_weight_sum').GetBinContent(1)
        xsec = h_xsec.GetBinContent(1) / h_xsec.GetEntries()
        h.Scale(xsec/weight_sum)

        #--------------------------------------------------
        # (2) Scale all histograms by bin width
        h.Scale(1., 'width')

        #--------------------------------------------------
        # (3) Scale for observable-specific factors
        if self.sqrts == 200:
            if observable_type == 'hadron':
                if observable == 'pt_ch_star':
                    sigma_inel = 42.
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                    h.Scale(1./sigma_inel)
                if observable == 'pt_pi0_phenix':
                    sigma_inel = 42. # Update
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                    h.Scale(1./sigma_inel)
            elif observable_type == 'hadron_correlations':
                if observable == "dihadron_star":
                    # Need to grab the pt trig range
                    # Example: "_pt_trig_8_15_...." -> split -> ['', 'pt', 'trig', '8', '15', '....']
                    # So we grab index 3 and 4
                    split_suffix = self.suffix.split("_")
                    pt_trig_min, pt_trig_max = float(split_suffix[3]), float(split_suffix[4])

                    # And then the trigger hist
                    hname = f"h_{observable_type}_dihadron_star_Ntrig_{centrality}"
                    h_ntrig = self.input_file.Get(hname)
                    h_ntrig.SetDirectory(0)

                    # Finally, scale by the n_trig
                    n_trig = h_ntrig.GetBinContent(h_ntrig.GetXaxis().FindBin((pt_trig_max + pt_trig_min) / 2)) # Note that Njets histogram should also be scaled by xsec/n_events
                    if n_trig > 0.:
                        h.Scale(1./(n_trig * (xsec/weight_sum)))
                    else:
                        print('WARNING: N_trig for dihadron_star == 0')
            elif observable_type == 'inclusive_chjet':
                if observable == 'pt_star':
                    h.Scale(1./(2*self.eta_cut))
            elif observable_type == 'semi_inclusive_chjet':
                # Note that n_trig histograms should also be scaled by xsec/n_events
                if observable in ['IAA_star', 'dphi_star']:
                    h.Scale(1./(xsec/weight_sum))

        if self.sqrts == 2760:
            if observable_type == 'hadron':
                if observable in ['pt_ch_alice', 'pt_pi_alice', 'pt_pi0_alice']:
                    sigma_inel = 61.8
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./sigma_inel)
                elif observable == 'pt_ch_atlas':
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                elif observable == 'pt_ch_cms':
                    sigma_inel = 64.
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                    h.Scale(1./sigma_inel)
            elif observable_type == 'inclusive_jet':
                if observable == 'pt_alice':
                    sigma_inel = 62.1
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./sigma_inel)
                if observable == 'pt_atlas':
                    h.Scale(1./(2*self.y_cut))
                    h.Scale(1.e6) # convert to nb
                if observable == 'pt_cms':
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1.e6) # convert to nb
                if observable in ['Dz_atlas', 'Dpt_atlas', 'Dz_cms', 'Dpt_cms']:
                    if observable in ['Dz_atlas', 'Dpt_atlas']:
                        hname = f'h_{observable_type}_Dz_atlas{self.suffix}{collection_label}_Njets_{centrality}{pt_suffix}'
                    elif observable in ['Dz_cms', 'Dpt_cms']:
                        hname = f'h_{observable_type}_Dz_cms{self.suffix}{collection_label}_Njets_{centrality}{pt_suffix}'
                    h_njets = self.input_file.Get(hname)
                    h_njets.SetDirectory(0)
                    n_jets = h_njets.GetBinContent(1) # Note that Njets histogram should also be scaled by xsec/n_events
                    if n_jets > 0.:
                        h.Scale(1./(n_jets * (xsec/weight_sum)))
                    else:
                        print('WARNING: N_jets = 0')
            elif observable_type == 'inclusive_chjet':
                if observable == 'pt_alice':
                    h.Scale(1./(2*self.eta_cut))
            elif observable_type == 'semi_inclusive_chjet':
                # Note that n_trig histograms should also be scaled by xsec/n_events
                if observable in ['IAA_alice', 'dphi_alice']:
                    h.Scale(1./(xsec/weight_sum))

        if self.sqrts == 5020:
            if observable_type == 'hadron':
                if observable in ['pt_ch_alice', 'pt_pi_alice']:
                    sigma_inel = 67.6
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./sigma_inel)
                elif observable == 'pt_ch_cms':
                    sigma = 70.
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1./(2*np.pi))
                    h.Scale(1./sigma)
            elif observable_type == 'inclusive_jet':
                if observable == 'pt_alice':
                    h.Scale(1./(2*self.eta_cut))
                if observable in ['pt_atlas', 'pt_y_atlas']:
                    h.Scale(1.e6) # convert to nb
                    if observable == 'pt_atlas':
                        h.Scale(1./(2*self.y_cut))
                if observable == 'pt_cms':
                    h.Scale(1./(2*self.eta_cut))
                    h.Scale(1.e6) # convert to nb
                if observable in ['Dz_atlas', 'Dpt_atlas']:
                    hname = f'h_{observable_type}_Dz_atlas{self.suffix}{collection_label}_Njets_{centrality}{pt_suffix}'

                    h_njets = self.input_file.Get(hname)
                    h_njets.SetDirectory(0)
                    n_jets = h_njets.GetBinContent(1) # Note that Njets histogram should also be scaled by xsec/n_events
                    if n_jets > 0.:
                        h.Scale(1./(n_jets * (xsec/weight_sum)))
                    else:
                        print('WARNING: N_jets = 0')

        # Scale by bin-dependent factor (approximation for pp QA)
        if self.scale_by:
            nBins = h.GetNbinsX()
            for bin in range(1, nBins+1):
                h_x = h.GetBinCenter(bin)
                h_y = h.GetBinContent(bin)
                h_error = h.GetBinError(bin)
                if self.scale_by == '1/pt':
                    scale_factor = 1./h_x
                h.SetBinContent(bin, h_y*scale_factor)
                h.SetBinError(bin, h_error*scale_factor)

        #--------------------------------------------------
        # (4) Self-normalize histogram if appropriate
        if self_normalize:
            if observable in ['zg_alice', 'tg_alice']:
                min_bin = 0
            else:
                min_bin = 1
            nbins = h.GetNbinsX()
            integral = h.Integral(min_bin, nbins, 'width')
            if integral > 0.:
                h.Scale(1./integral)
            else:
                print('WARNING: integral for self-normalization is 0')

    #-------------------------------------------------------------------------------------------
    # Perform any additional manipulations on scaled histograms
    #-------------------------------------------------------------------------------------------
    def post_process_histogram(self, observable_type, observable, block, centrality, centrality_index: int, collection_label=''):
        if 'v2' in observable and self.is_AA:
            # hadron v2
            h = self.observable_settings[f'jetscape_distribution{collection_label}']
            if h:
                h_num_name = f'h_{observable_type}_{observable}_{centrality}'
                h_num_name_holes = f'h_{observable_type}_{observable}_holes_{centrality}'
                h_denom_name = f'h_{observable_type}_{observable}_denom_{centrality}'
                h_denom_name_holes = f'h_{observable_type}_{observable}_holes_denom_{centrality}'

                self.observable_settings[f'jetscape_distribution'] = self.input_file.Get(h_num_name).Clone()
                self.observable_settings[f'jetscape_distribution_unsubtracted'] = self.input_file.Get(h_num_name).Clone()

                for i in range(0,self.input_file.Get(h_num_name).GetNbinsX()):
                    h_num_i = self.input_file.Get(h_num_name).GetBinContent(i)
                    h_num_holes_i = self.input_file.Get(h_num_name_holes).GetBinContent(i)
                    h_denom_i = self.input_file.Get(h_denom_name).GetBinContent(i)
                    h_denom_holes_i = self.input_file.Get(h_denom_name_holes).GetBinContent(i)
                    if h_denom_i == 0.0:
                        h_denom_i = -99.0
                        h_denom_holes_i =0.0
                    self.observable_settings[f'jetscape_distribution'].SetBinContent(i, (h_num_i - h_num_holes_i)/(h_denom_i - h_denom_holes_i) )
                    self.observable_settings[f'jetscape_distribution_unsubtracted'].SetBinContent(i,h_num_i/h_denom_i)

                self.observable_settings[f'jetscape_distribution'].Print()
                self.observable_settings[f'jetscape_distribution_unsubtracted'].Print()

                self.observable_settings[f'jetscape_distribution'].SetName(f'jetscape_distribution_{observable_type}_{observable}_{centrality}')
                self.observable_settings[f'jetscape_distribution_unsubtracted'].SetName(f'jetscape_distribution_unsubtracted_{observable_type}_{observable}_{centrality}')

        # For the ATLAS rapidity-dependence, we need to divide the histograms by their first bin (|y|<0.3) to form a double ratio
        if observable == 'pt_y_atlas':

            # Get y=0.3 entry, and create a histogram to divide by (so that uncertainties are propagated)
            h = self.observable_settings[f'jetscape_distribution{collection_label}']
            if h:
                denominator = h.GetBinContent(1)
                denominator_uncertainty = h.GetBinError(1)
                h_denominator = h.Clone()
                h_denominator.SetName(f'{h_denominator.GetName()}_denominator')
                n = h_denominator.GetNbinsX()
                for i in range(n):
                    h_denominator.SetBinContent(i+1, denominator)
                    h_denominator.SetBinError(i+1, denominator_uncertainty)

                # Divide histograms
                h.Divide(h_denominator)

                # Rebin to remove first bin (|y|<0.3)
                bins = np.array(h.GetXaxis().GetXbins())
                hname = h.GetName()
                self.observable_settings[f'jetscape_distribution{collection_label}'] = h.Rebin(n-1, hname, bins[1:])

        if observable_type == "hadron_correlations" and observable == "dihadron_star":
            # In the case of pp, just look at the deltaPhi correlations themselves
            # In any case, we don't have data to compare against.
            # If skip_pp is turned off, then we can check the yields in pp.
            if not self.skip_pp:
                # Grab the delta phi correlation
                h = self.observable_settings[f'jetscape_distribution']
                # Extract the yield
                yield_lower_range_value, yield_upper_range_value = block["yield_range"]
                # First, near side
                _temp_error = ctypes.c_double(0)
                near_side_yield = h.IntegralAndError(
                    h.GetXaxis().FindBin(yield_lower_range_value), h.GetXaxis().FindBin(yield_upper_range_value),
                    _temp_error,
                )
                near_side_yield_error = _temp_error.value
                # Then the away side
                _temp_error = ctypes.c_double(0)
                away_side_yield = h.IntegralAndError(
                    # Both are np.pi + because the yield_lower_range_value is negative
                    h.GetXaxis().FindBin(np.pi + yield_lower_range_value), h.GetXaxis().FindBin(np.pi + yield_upper_range_value),
                    _temp_error,
                )
                away_side_yield_error = _temp_error.value

                # Encode the values into a single TH1F. We put it into a TH1F instead of a TGraph because
                # the framework seems to require that the jetscape data is in a histogram.
                yield_bins = np.array(block["yield_bins"])
                n = len(yield_bins) - 1
                h_yield = ROOT.TH1F(f"{h.GetName()}_yield", f"{h.GetName()}_yields", n, yield_bins)

                # Fill in extracted values
                yield_values = np.zeros(n)
                yield_errors = np.zeros(n)
                yield_values[centrality_index] = near_side_yield
                yield_errors[centrality_index] = near_side_yield_error
                # round is just to ensure we have an int
                yield_values[centrality_index + round(n / 2)] = away_side_yield
                yield_errors[centrality_index + round(n / 2)] = away_side_yield_error

                # start at 1 to account for ROOT bin indexing
                for i, (yv, ye) in enumerate(zip(yield_values, yield_errors), start=1):
                    h_yield.SetBinContent(i, yv)
                    h_yield.SetBinError(i, ye)

                # Now, put it together and store the result
                self.observable_settings["jetscape_distribution"] = h_yield

    #-------------------------------------------------------------------------------------------
    # Plot distributions in upper panel, and ratio in lower panel
    #-------------------------------------------------------------------------------------------
    def plot_observable(self, observable_type, observable, centrality, pt_suffix='', logy = False):

        label = f'{observable_type}_{observable}_{self.sqrts}_{centrality}_{self.suffix}_{pt_suffix}'

        if 'v2' in observable:
            # for hadron v2
            if self.observable_settings[f'jetscape_distribution']:
                self.plot_distribution_and_ratio(observable_type, observable, centrality, label, pt_suffix=pt_suffix, logy=logy)
            return

        # If AA: Plot PbPb/pp ratio, and comparison to data
        # If pp: Plot distribution, and ratio to data
        if self.is_AA:
            self.plot_RAA(observable_type, observable, centrality, label, pt_suffix=pt_suffix, logy=logy)
        else:
            if self.observable_settings[f'jetscape_distribution']:
                self.plot_distribution_and_ratio(observable_type, observable, centrality, label, pt_suffix=pt_suffix, logy=logy)

    #-------------------------------------------------------------------------------------------
    # Plot distributions in upper panel, and ratio in lower panel
    # We plot all available modes of hole subtraction:
    #   For hadron histograms (self.subtract_holes = True):
    #       - self.observable_settings['jetscape_distribution']
    #       - self.observable_settings['jetscape_distribution_unsubtracted']
    #   For jet histograms:
    #       for jet_collection_label in self.jet_collection_labels:
    #           self.observable_settings[f'jetscape_distribution{jet_collection_label}']
    #-------------------------------------------------------------------------------------------
    def plot_RAA(self, observable_type, observable, centrality, label, pt_suffix='', logy = False):

        # Assemble the list of hole subtraction variations
        keys_to_plot = [key for key in self.observable_settings.keys() if 'jetscape_distribution' in key and 'holes' not in key]
        self.jetscape_legend_label = {}
        self.jetscape_legend_label['jetscape_distribution'] = 'JETSCAPE'
        self.jetscape_legend_label['jetscape_distribution_unsubtracted'] = 'JETSCAPE (unsubtracted)'
        self.jetscape_legend_label['jetscape_distribution_shower_recoil'] = 'JETSCAPE (shower+recoil)'
        self.jetscape_legend_label['jetscape_distribution_shower_recoil_unsubtracted'] = 'JETSCAPE (shower+recoil, unsubtracted)'
        self.jetscape_legend_label['jetscape_distribution_negative_recombiner'] = 'JETSCAPE (negative recombiner)'
        self.jetscape_legend_label['jetscape_distribution_constituent_subtraction'] = 'JETSCAPE (CS)'

        if not self.observable_settings[keys_to_plot[0]]:
            print(f'WARNING: skipping {observable} {label} {centrality} since data is missing')
            return

        # Get the pp reference histogram and form the RAA ratios
        if not self.skip_AA_ratio:
            h_pp = self.pp_ref_file.Get(f'jetscape_distribution_{label}')
            for key in keys_to_plot:
                if self.observable_settings[key] and h_pp:
                    self.observable_settings[key].Divide(h_pp)

        c = ROOT.TCanvas('c', 'c', 600, 450)
        c.SetRightMargin(0.05);
        c.SetLeftMargin(0.15);
        c.SetTopMargin(0.05);
        c.SetBottomMargin(0.17);
        c.cd()

        legend = ROOT.TLegend(0.4,0.55,0.75,0.81)
        self.plot_utils.setup_legend(legend, 0.045, sep=-0.1)

        if not self.skip_AA_ratio:
            self.bins = np.array(h_pp.GetXaxis().GetXbins())
        else:
            self.bins = np.array(self.observable_settings[keys_to_plot[0]].GetXaxis().GetXbins())

        myBlankHisto = ROOT.TH1F('myBlankHisto','Blank Histogram', 1, self.bins[0], self.bins[-1])
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle(self.xtitle)
        myBlankHisto.SetYTitle(self.ytitle)
        myBlankHisto.SetMaximum(self.y_max)
        myBlankHisto.SetMinimum(self.y_min)
        myBlankHisto.GetXaxis().SetTitleSize(0.07)
        myBlankHisto.GetYaxis().SetTitleSize(0.07)
        myBlankHisto.GetXaxis().SetTitleOffset(1.)
        myBlankHisto.GetYaxis().SetTitleOffset(1.)
        myBlankHisto.Draw('E')

        # Draw data
        if self.observable_settings['data_distribution']:
            self.output_dict[f'data_distribution_{label}'] = self.observable_settings['data_distribution']
            self.observable_settings['data_distribution'].SetName(f'data_distribution_{label}')
            self.observable_settings['data_distribution'].SetMarkerSize(self.marker_size)
            self.observable_settings['data_distribution'].SetMarkerStyle(self.data_marker)
            self.observable_settings['data_distribution'].SetMarkerColor(self.data_color)
            self.observable_settings['data_distribution'].SetLineStyle(self.line_style)
            self.observable_settings['data_distribution'].SetLineWidth(self.line_width)
            self.observable_settings['data_distribution'].SetLineColor(self.data_color)
            self.observable_settings['data_distribution'].Draw('PE Z same')
            legend.AddEntry(self.observable_settings['data_distribution'], 'Data', 'PE')

        # Draw JETSCAPE
        for i,key in enumerate(keys_to_plot):
            self.output_dict[f'{key}_{label}'] = self.observable_settings[key]
            if self.observable_settings[key]:
                if self.observable_settings[key].GetNbinsX() > 1:
                    self.observable_settings[key].SetFillColor(self.jetscape_color[i])
                    self.observable_settings[key].SetFillColorAlpha(self.jetscape_color[i], self.jetscape_alpha[i])
                    self.observable_settings[key].SetFillStyle(self.jetscape_fillstyle[i])
                    self.observable_settings[key].SetMarkerSize(0.)
                    self.observable_settings[key].SetMarkerStyle(0)
                    self.observable_settings[key].SetLineWidth(0)
                    self.observable_settings[key].DrawCopy('E3 same')
                elif self.observable_settings[key].GetNbinsX() == 1:
                    self.observable_settings[key].SetMarkerSize(self.marker_size)
                    self.observable_settings[key].SetMarkerStyle(self.data_marker+1)
                    self.observable_settings[key].SetMarkerColor(self.jetscape_color[i])
                    self.observable_settings[key].SetLineStyle(self.line_style)
                    self.observable_settings[key].SetLineWidth(self.line_width)
                    self.observable_settings[key].SetLineColor(self.jetscape_color[i])
                    self.observable_settings[key].DrawCopy('PE same')
                legend.AddEntry(self.observable_settings[key], self.jetscape_legend_label[key], 'f')

        legend.Draw()

        if self.y_max > 1 and not self.skip_AA_ratio:
            line = ROOT.TLine(self.bins[0], 1, self.bins[-1], 1)
            line.SetLineColor(920+2)
            line.SetLineStyle(2)
            line.SetLineWidth(2)
            line.Draw()

        text_latex = ROOT.TLatex()
        text_latex.SetNDC()

        x = 0.18
        text_latex.SetTextSize(0.065)
        text = f'#bf{{{observable_type}_{observable}}} #sqrt{{#it{{s}}}} = {self.sqrts/1000.} TeV'
        text_latex.DrawLatex(x, 0.88, text)
        text = f'{centrality} {self.suffix} {pt_suffix}'
        text_latex.DrawLatex(x, 0.8, text)

        hname = f'h_{observable_type}_{observable}{self.suffix}_{centrality}{pt_suffix}'
        c.SaveAs(os.path.join(self.output_dir, f'{hname}{self.file_format}'))
        c.Close()

    #-------------------------------------------------------------------------------------------
    # Plot distributions in upper panel, and ratio in lower panel
    #-------------------------------------------------------------------------------------------
    def plot_distribution_and_ratio(self, observable_type, observable, centrality, label, pt_suffix='', logy = False):

        c = ROOT.TCanvas('c', 'c', 600, 650)
        c.Draw()
        c.cd()

        # Distribution
        pad2_dy = 0.45
        pad1 = ROOT.TPad('myPad', 'The pad',0,pad2_dy,1,1)
        pad1.SetLeftMargin(0.2)
        pad1.SetTopMargin(0.08)
        pad1.SetRightMargin(0.04)
        pad1.SetBottomMargin(0.)
        pad1.SetTicks(0,1)
        pad1.Draw()
        if self.logy:
            pad1.SetLogy()
        pad1.cd()

        legend = ROOT.TLegend(0.58,0.6,0.75,0.75)
        self.plot_utils.setup_legend(legend, 0.055, sep=-0.1)

        legend_ratio = ROOT.TLegend(0.25,0.8,0.45,1.07)
        self.plot_utils.setup_legend(legend_ratio, 0.07, sep=-0.1)

        self.bins = np.array(self.observable_settings[f'jetscape_distribution'].GetXaxis().GetXbins())
        myBlankHisto = ROOT.TH1F('myBlankHisto','Blank Histogram', 1, self.bins[0], self.bins[-1])
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle(self.xtitle)
        myBlankHisto.SetYTitle(self.ytitle)
        myBlankHisto.SetMaximum(self.y_max)
        myBlankHisto.SetMinimum(self.y_min) # Don't draw 0 on top panel
        myBlankHisto.GetYaxis().SetTitleSize(0.08)
        myBlankHisto.GetYaxis().SetTitleOffset(1.1)
        myBlankHisto.GetYaxis().SetLabelSize(0.06)
        myBlankHisto.Draw('E')

        # Ratio
        c.cd()
        pad2 = ROOT.TPad('pad2', 'pad2', 0, 0.02, 1, pad2_dy)
        pad2.SetTopMargin(0)
        pad2.SetBottomMargin(0.21)
        pad2.SetLeftMargin(0.2)
        pad2.SetRightMargin(0.04)
        pad2.SetTicks(0,1)
        pad2.Draw()
        pad2.cd()

        myBlankHisto2 = myBlankHisto.Clone('myBlankHisto_C')
        myBlankHisto2.SetYTitle('Ratio')
        myBlankHisto2.SetXTitle(self.xtitle)
        myBlankHisto2.GetXaxis().SetTitleSize(26)
        myBlankHisto2.GetXaxis().SetTitleFont(43)
        myBlankHisto2.GetXaxis().SetTitleOffset(2.3)
        myBlankHisto2.GetXaxis().SetLabelFont(43)
        myBlankHisto2.GetXaxis().SetLabelSize(22)
        myBlankHisto2.GetYaxis().SetTitleSize(28)
        myBlankHisto2.GetYaxis().SetTitleFont(43)
        myBlankHisto2.GetYaxis().SetTitleOffset(2.)
        myBlankHisto2.GetYaxis().SetLabelFont(43)
        myBlankHisto2.GetYaxis().SetLabelSize(20)
        myBlankHisto2.GetYaxis().SetNdivisions(505)
        myBlankHisto2.GetYaxis().SetRangeUser(self.y_ratio_min, self.y_ratio_max)
        myBlankHisto2.Draw('')

        # Draw distribution
        pad1.cd()
        if self.observable_settings['data_distribution']:
            self.output_dict[f'data_distribution_{label}'] = self.observable_settings['data_distribution']
            self.observable_settings['data_distribution'].SetName(f'data_distribution_{label}')
            self.observable_settings['data_distribution'].SetMarkerSize(self.marker_size)
            self.observable_settings['data_distribution'].SetMarkerStyle(self.data_marker)
            self.observable_settings['data_distribution'].SetMarkerColor(self.data_color)
            self.observable_settings['data_distribution'].SetLineStyle(self.line_style)
            self.observable_settings['data_distribution'].SetLineWidth(self.line_width)
            self.observable_settings['data_distribution'].SetLineColor(self.data_color)
            self.observable_settings['data_distribution'].Draw('PE Z same')
            legend.AddEntry(self.observable_settings['data_distribution'], 'Data', 'PE')

        # Draw JETSCAPE
        self.output_dict[f'jetscape_distribution_{label}'] =  self.observable_settings[f'jetscape_distribution']
        if self.observable_settings[f'jetscape_distribution'].GetNbinsX() > 1:
            self.observable_settings[f'jetscape_distribution'].SetFillColor(self.jetscape_color[0])
            self.observable_settings[f'jetscape_distribution'].SetFillColorAlpha(self.jetscape_color[0], self.jetscape_alpha[0])
            self.observable_settings[f'jetscape_distribution'].SetFillStyle(1001)
            self.observable_settings[f'jetscape_distribution'].SetMarkerSize(0.)
            self.observable_settings[f'jetscape_distribution'].SetMarkerStyle(0)
            self.observable_settings[f'jetscape_distribution'].SetLineWidth(0)
            self.observable_settings[f'jetscape_distribution'].DrawCopy('E3 same')
        elif self.observable_settings[f'jetscape_distribution'].GetNbinsX() == 1:
            self.observable_settings[f'jetscape_distribution'].SetMarkerSize(self.marker_size)
            self.observable_settings[f'jetscape_distribution'].SetMarkerStyle(self.data_marker+1)
            self.observable_settings[f'jetscape_distribution'].SetMarkerColor(self.jetscape_color[0])
            self.observable_settings[f'jetscape_distribution'].SetLineStyle(self.line_style)
            self.observable_settings[f'jetscape_distribution'].SetLineWidth(self.line_width)
            self.observable_settings[f'jetscape_distribution'].SetLineColor(self.jetscape_color[0])
            self.observable_settings[f'jetscape_distribution'].DrawCopy('PE same')
        legend.AddEntry(self.observable_settings[f'jetscape_distribution'], 'JETSCAPE', 'f')

        legend.Draw()

        # Draw ratio
        pad2.cd()
        if self.observable_settings['data_distribution']:
            data_ratio = self.plot_utils.divide_tgraph_by_tgraph(self.observable_settings['data_distribution'],
                                                                 self.observable_settings['data_distribution'])
            data_ratio.Draw('PE Z same')
            self.output_dict[f'data_ratio_{label}'] = data_ratio
        if self.observable_settings['ratio']:
            self.output_dict[f'ratio_{label}'] = self.observable_settings['ratio']
            if self.observable_settings['ratio'].GetN() > 1:
                self.observable_settings['ratio'].SetFillColor(self.jetscape_color[0])
                self.observable_settings['ratio'].SetFillColorAlpha(self.jetscape_color[0], self.jetscape_alpha[0])
                self.observable_settings['ratio'].SetFillStyle(1001)
                self.observable_settings['ratio'].SetMarkerSize(0.)
                self.observable_settings['ratio'].SetMarkerStyle(0)
                self.observable_settings['ratio'].SetLineWidth(0)
                self.observable_settings['ratio'].Draw('E3 same')
            elif self.observable_settings['ratio'].GetN() == 1:
                self.observable_settings['ratio'].SetMarkerSize(self.marker_size)
                self.observable_settings['ratio'].SetMarkerStyle(self.data_marker+1)
                self.observable_settings['ratio'].SetMarkerColor(self.jetscape_color[0])
                self.observable_settings['ratio'].SetLineStyle(self.line_style)
                self.observable_settings['ratio'].SetLineWidth(self.line_width)
                self.observable_settings['ratio'].SetLineColor(self.jetscape_color[0])
                self.observable_settings['ratio'].Draw('PE same')

        if self.observable_settings['ratio']:
            legend_ratio.AddEntry(self.observable_settings['ratio'], 'JETSCAPE/Data', 'f')
        if self.observable_settings['data_distribution']:
            legend_ratio.AddEntry(data_ratio, 'Data uncertainties', 'PE')
        legend_ratio.Draw()

        line = ROOT.TLine(self.bins[0], 1, self.bins[-1], 1)
        line.SetLineColor(920+2)
        line.SetLineStyle(2)
        line.SetLineWidth(2)
        line.Draw()

        pad1.cd()
        text_latex = ROOT.TLatex()
        text_latex.SetNDC()

        x = 0.25
        text_latex.SetTextSize(0.065)
        text = f'#bf{{{observable_type}_{observable}}} #sqrt{{#it{{s}}}} = {self.sqrts/1000.} TeV'
        text_latex.DrawLatex(x, 0.83, text)
        text = f'{centrality} {self.suffix} {pt_suffix}'
        text_latex.DrawLatex(x, 0.73, text)

        if self.skip_pp:
            text = 'skip data plot -- no pp data in HEPData'
            text_latex.DrawLatex(x, 0.53, text)

        pad2.cd()
        if self.skip_pp_ratio:
            text = 'skip ratio plot -- pp/AA binning mismatch'
            text_latex.DrawLatex(x, 0.93, text)

        c.SaveAs(os.path.join(self.output_dir, f'{self.hname}{self.file_format}'))
        c.Close()

    #-------------------------------------------------------------------------------------------
    # Plot event QA
    #-------------------------------------------------------------------------------------------
    def plot_event_qa(self):

        if self.is_AA:

            for centrality in self.observable_centrality_list:

                # Only plot those centralities that exist
                if np.isclose(self.input_file.Get(f'h_centrality_generated').Integral(centrality[0]+1, centrality[1]), 0):
                    continue
                print(centrality)

                # Crosscheck that pt-hat weighting is satisfactory
                # Make sure that the large-weight, low-pt-hat range has sufficient statistics to avoid normalization fluctuations
                sum_weights = self.input_file.Get(f'h_weight_sum_{centrality}').GetBinContent(1)
                h_pt_hat = self.input_file.Get('h_pt_hat')
                h_pt_hat_weighted = self.input_file.Get('h_pt_hat_weighted')

                # Normalize by sum of weights, i.e. 1/sigma_pt_hat * dsigma/dpt_hat
                h_pt_hat.Scale(1./sum_weights, 'width')
                h_pt_hat_weighted.Scale(1./sum_weights, 'width')

                # Compute normalization uncertainty (binned approximation)
                h_weights = self.input_file.Get(f'h_weights_{centrality}')
                sum_weights_integral = 0
                normalization_uncertainty = 0
                for i in range(1, h_weights.GetNbinsX()+1):
                    sum_weights_integral += h_weights.GetBinCenter(i)*h_weights.GetBinContent(i)
                    normalization_uncertainty = np.sqrt( np.square(normalization_uncertainty) + h_weights.GetBinContent(i)*np.square(h_weights.GetBinCenter(i)) )
                print(f'sum_weights {centrality}: {sum_weights}')
                print(f'sum_weights_integral {centrality}: {sum_weights_integral}')
                print(f'normalization_uncertainty {centrality}: {100*normalization_uncertainty/sum_weights_integral} %')

                # Also compute overall pt-hat cross-section uncertainty
                h_xsec = self.input_file.Get(f'h_xsec_{centrality}')
                xsec = h_xsec.GetBinContent(1)
                xsec_error = self.input_file.Get(f'h_xsec_error_{centrality}').GetBinContent(1)
                print(f'xsec {centrality}: {xsec/h_xsec.GetEntries()}')
                print(f'xsec_uncertainty {centrality}: {100*xsec_error/xsec}')
                print()

        else:

            # Crosscheck that pt-hat weighting is satisfactory
            # Make sure that the large-weight, low-pt-hat range has sufficient statistics to avoid normalization fluctuations
            sum_weights = self.input_file.Get('h_weight_sum').GetBinContent(1)
            h_pt_hat = self.input_file.Get('h_pt_hat')
            h_pt_hat_weighted = self.input_file.Get('h_pt_hat_weighted')

            # Normalize by sum of weights, i.e. 1/sigma_pt_hat * dsigma/dpt_hat
            h_pt_hat.Scale(1./sum_weights, 'width')
            h_pt_hat_weighted.Scale(1./sum_weights, 'width')

            # Compute normalization uncertainty (binned approximation)
            h_weights = self.input_file.Get('h_weights')
            sum_weights_integral = 0
            normalization_uncertainty = 0
            for i in range(1, h_weights.GetNbinsX()+1):
                sum_weights_integral += h_weights.GetBinCenter(i)*h_weights.GetBinContent(i)
                normalization_uncertainty = np.sqrt( np.square(normalization_uncertainty) + h_weights.GetBinContent(i)*np.square(h_weights.GetBinCenter(i)) )
            print(f'sum_weights: {sum_weights}')
            print(f'sum_weights_integral: {sum_weights_integral}')
            print(f'normalization_uncertainty: {100*normalization_uncertainty/sum_weights_integral} %')

            # Also compute overall pt-hat cross-section uncertainty
            h_xsec = self.input_file.Get('h_xsec')
            xsec = h_xsec.GetBinContent(1)
            xsec_error = self.input_file.Get('h_xsec_error').GetBinContent(1)
            print(f'xsec: {xsec/h_xsec.GetEntries()}')
            print(f'xsec_uncertainty: {100*xsec_error/xsec}')

        c = ROOT.TCanvas('c', 'c', 600, 650)
        c.Draw()
        c.cd()

        # Distribution
        pad2_dy = 0.45
        pad1 = ROOT.TPad('myPad', 'The pad',0,pad2_dy,1,1)
        pad1.SetLeftMargin(0.2)
        pad1.SetTopMargin(0.08)
        pad1.SetRightMargin(0.04)
        pad1.SetBottomMargin(0.)
        pad1.SetTicks(0,1)
        pad1.Draw()
        pad1.SetLogy()
        pad1.SetLogx()
        pad1.cd()

        legend = ROOT.TLegend(0.58,0.65,0.75,0.9)
        self.plot_utils.setup_legend(legend, 0.055, sep=-0.1, title=f'#sqrt{{s_{{NN}}}} = {self.sqrts} GeV')

        self.bins = np.array(h_pt_hat.GetXaxis().GetXbins())
        myBlankHisto = ROOT.TH1F('myBlankHisto','Blank Histogram', 1, self.bins[0], self.bins[-1])
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle('#hat{p}_{T} (GeV/#it{c})')
        myBlankHisto.SetYTitle('#frac{1}{#sigma_{#hat{p}_{T}}} #frac{d#sigma}{d#hat{p}_{T}}')
        myBlankHisto.SetMaximum(np.power(10,5+(self.power-3)))
        myBlankHisto.SetMinimum(2e-15) # Don't draw 0 on top panel
        myBlankHisto.GetYaxis().SetTitleSize(0.08)
        myBlankHisto.GetYaxis().SetTitleOffset(1.1)
        myBlankHisto.GetYaxis().SetLabelSize(0.06)
        myBlankHisto.Draw('E')

        # Ratio
        c.cd()
        pad2 = ROOT.TPad('pad2', 'pad2', 0, 0.02, 1, pad2_dy)
        pad2.SetTopMargin(0)
        pad2.SetBottomMargin(0.21)
        pad2.SetLeftMargin(0.2)
        pad2.SetRightMargin(0.04)
        pad2.SetTicks(0,1)
        pad2.SetLogx()
        pad2.Draw()
        pad2.cd()

        myBlankHisto2 = myBlankHisto.Clone('myBlankHisto_C')
        myBlankHisto2.SetYTitle('#frac{weighted}{ideal}')
        myBlankHisto2.SetXTitle('#hat{p}_{T} (GeV/#it{c})')
        myBlankHisto2.GetXaxis().SetTitleSize(26)
        myBlankHisto2.GetXaxis().SetTitleFont(43)
        myBlankHisto2.GetXaxis().SetTitleOffset(2.3)
        myBlankHisto2.GetXaxis().SetLabelFont(43)
        myBlankHisto2.GetXaxis().SetLabelSize(22)
        myBlankHisto2.GetYaxis().SetTitleSize(28)
        myBlankHisto2.GetYaxis().SetTitleFont(43)
        myBlankHisto2.GetYaxis().SetTitleOffset(2.)
        myBlankHisto2.GetYaxis().SetLabelFont(43)
        myBlankHisto2.GetYaxis().SetLabelSize(20)
        myBlankHisto2.GetYaxis().SetNdivisions(505)
        myBlankHisto2.GetYaxis().SetRangeUser(0., 1.99)
        myBlankHisto2.Draw('')

        # Draw generated pt-hat
        pad1.cd()
        h_pt_hat.SetMarkerSize(self.marker_size)
        h_pt_hat.SetMarkerStyle(self.data_marker)
        h_pt_hat.SetMarkerColor(self.data_color)
        h_pt_hat.SetLineStyle(self.line_style)
        h_pt_hat.SetLineWidth(self.line_width)
        h_pt_hat.SetLineColor(self.data_color)
        h_pt_hat.Draw('PE Z same')
        legend.AddEntry(h_pt_hat, f'#alpha={self.power}, generated', 'PE')

        # Draw weighted pt-hat
        h_pt_hat_weighted.SetMarkerSize(self.marker_size)
        h_pt_hat_weighted.SetMarkerStyle(self.data_marker)
        h_pt_hat_weighted.SetMarkerColor(self.jetscape_color[0])
        h_pt_hat_weighted.SetLineStyle(self.line_style)
        h_pt_hat_weighted.SetLineWidth(self.line_width)
        h_pt_hat_weighted.SetLineColor(self.jetscape_color[0])
        h_pt_hat_weighted.Draw('PE Z same')
        legend.AddEntry(h_pt_hat_weighted, f'#alpha={self.power}, weighted', 'PE')

        # Draw generated pt-hat weighted by inverse of ideal weight (pthat/ptref)^alpha
        h_pt_hat_weighted_ideal = h_pt_hat.Clone()
        h_pt_hat_weighted_ideal.SetName('h_pt_hat_weighted_ideal')
        for i in range(1, h_pt_hat_weighted_ideal.GetNbinsX()+1):
            content = h_pt_hat_weighted_ideal.GetBinContent(i)
            low_edge = h_pt_hat_weighted_ideal.GetXaxis().GetBinLowEdge(i)
            up_edge = h_pt_hat_weighted_ideal.GetXaxis().GetBinUpEdge(i)
            ideal_weight = ( np.power(up_edge, self.power+1)-np.power(low_edge, self.power+1) ) / ( (up_edge-low_edge)*np.power(self.pt_ref, self.power)*(self.power+1) )
            #ideal_weight = np.power( h_pt_hat_weighted_ideal.GetXaxis().GetBinCenter(i) / self.pt_ref, self.power)
            h_pt_hat_weighted_ideal.SetBinContent(i, content / ideal_weight)
            h_pt_hat_weighted_ideal.SetBinError(i, 0.)
        h_pt_hat_weighted_ideal.SetLineColorAlpha(ROOT.kGreen-8, self.jetscape_alpha[0])
        h_pt_hat_weighted_ideal.SetLineWidth(3)
        h_pt_hat_weighted_ideal.Draw('L same')
        legend.AddEntry(h_pt_hat_weighted_ideal, f'#alpha={self.power}, ideal weight', 'L')

        legend.Draw()

        # Plot ratio of weighted to ideal weighted
        pad2.cd()
        h_pt_hat_weighted_ratio = h_pt_hat_weighted.Clone()
        h_pt_hat_weighted_ratio.SetName('h_pt_hat_weighted_ratio')
        h_pt_hat_weighted_ratio.Divide(h_pt_hat_weighted_ideal)
        h_pt_hat_weighted_ratio.SetMarkerSize(self.marker_size)
        h_pt_hat_weighted_ratio.SetMarkerStyle(self.data_marker)
        h_pt_hat_weighted_ratio.SetMarkerColor(self.jetscape_color[0])
        h_pt_hat_weighted_ratio.SetLineStyle(self.line_style)
        h_pt_hat_weighted_ratio.SetLineWidth(self.line_width)
        h_pt_hat_weighted_ratio.SetLineColor(self.jetscape_color[0])
        h_pt_hat_weighted_ratio.Draw('PE Z same')

        line = ROOT.TLine(self.bins[0], 1, self.bins[-1], 1)
        line.SetLineColor(920+2)
        line.SetLineStyle(2)
        line.SetLineWidth(2)
        line.Draw()

        pad1.cd()
        text_latex = ROOT.TLatex()
        text_latex.SetNDC()

        x = 0.25
        text_latex.SetTextSize(0.06)
        text = '#sigma_{{n_{{event}}}} = {:.2f}%'.format(100*normalization_uncertainty/sum_weights_integral)
        text_latex.DrawLatex(x, 0.83, text)
        text = '#sigma_{{xsec}} = {:.2f}%'.format(100*xsec_error/xsec)
        text_latex.DrawLatex(x, 0.76, text)

        c.SaveAs(os.path.join(self.output_dir, f'pt_hat{self.file_format}'))
        c.Close()

    # ---------------------------------------------------------------
    # Save all ROOT histograms to file
    # ---------------------------------------------------------------
    def write_output_objects(self):
        print()

        # Save histograms to ROOT file
        output_ROOT_filename = os.path.join(self.output_dir, 'final_results.root')
        f_ROOT = ROOT.TFile(output_ROOT_filename, 'recreate')
        f_ROOT.cd()
        for key,val in self.output_dict.items():
            if val:
                val.SetName(key)
                val.Write()
                if isinstance(val, (ROOT.TH1)):
                    val.SetDirectory(0)
                    del val
        f_ROOT.Close()

        # Also save JETSCAPE AA/pp ratios as numpy arrays to HDF5 file
        if self.is_AA:
            with uproot.open(output_ROOT_filename) as uproot_file:
                output_HDF5_filename = os.path.join(self.output_dir, 'final_results.h5')
                with h5py.File(output_HDF5_filename, 'w') as hf:

                    keys = [key.split(';',1)[0] for key in uproot_file.keys()]
                    for key in keys:
                        if 'jetscape' in key:

                            bin_edges = uproot_file[key].axis().edges()
                            values = uproot_file[key].values()
                            errors = uproot_file[key].errors()

                            hf.create_dataset(f'{key}_bin_edges', data=bin_edges)
                            hf.create_dataset(f'{key}_values', data=values)
                            hf.create_dataset(f'{key}_errors', data=errors)

    #-------------------------------------------------------------------------------------------
    # Generate pptx of one plot per slide, for convenience
    #-------------------------------------------------------------------------------------------
    def generate_pptx(self):
        # Delayed import since this isn't on available on many systems
        # For those that need can install it (ie. systems where powerpoint can run) you can use:
        #   pip install python-pptx
        import pptx

        # Create a blank presentation
        p = pptx.Presentation()

        # Set slide layouts
        title_slide_layout = p.slide_layouts[0]
        blank_slide_layout = p.slide_layouts[6]

        # Make a title slide
        slide = p.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f'QA for {self.sqrts/1000.} TeV'
        author = slide.placeholders[1]
        author.text = 'STAT WG'

        # Loop through all output files and plot
        files = [f for f in os.listdir(self.output_dir) if f.endswith(self.file_format)]
        for file in sorted(files):
            img = os.path.join(self.output_dir, file)
            slide = p.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img, left=pptx.util.Inches(2.),
                                          top=pptx.util.Inches(1.),
                                          width=pptx.util.Inches(5.))

        p.save(os.path.join(self.output_dir, 'results.pptx'))

#-------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------
#-------------------------------------------------------------------------------------------
if __name__ == '__main__':
    print('Executing plot_results_STAT.py...')
    print('')

    # Define arguments
    parser = argparse.ArgumentParser(description='Plot JETSCAPE events')
    parser.add_argument(
        '-c',
        '--configFile',
        action='store',
        type=str,
        metavar='configFile',
        default='config/TG3.yaml',
        help='Config file'
    )
    parser.add_argument(
        '-i',
        '--inputFile',
        action='store',
        type=str,
        metavar='inputFile',
        default='pp_5020_plot/histograms_5020_merged.root',
        help='Input file'
    )
    parser.add_argument(
        '-o',
        '--outputDir',
        action='store',
        type=str,
        metavar='outputDir',
        default='',
        help='Output directory for output to be written to',
    )
    parser.add_argument(
        '-r',
        '--refFile',
        action='store',
        type=str,
        metavar='refFile',
        default='final_results.root',
        help='pp reference file'
    )

    # Parse the arguments
    args = parser.parse_args()

    # If invalid configFile is given, exit
    if not os.path.exists(args.configFile):
        print('File "{0}" does not exist! Exiting!'.format(args.configFile))
        sys.exit(0)

    # If invalid inputDir is given, exit
    if not os.path.exists(args.inputFile):
        print('File "{0}" does not exist! Exiting!'.format(args.inputFile))
        sys.exit(0)

    analysis = PlotResults(config_file=args.configFile, input_file=args.inputFile, output_dir=args.outputDir, pp_ref_file=args.refFile)
    analysis.plot_results()
